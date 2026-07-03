import sys
from pathlib import Path
from pptx import Presentation
from datetime import datetime

def main():
    pptx_path = Path("/Users/Abir/Desktop/innovation project/scale up/Three‑Engine Scale‑Up Canvas-UrbanGlowGrid.pptx")
    if not pptx_path.exists():
        print(f"PowerPoint file not found: {pptx_path}")
        sys.exit(1)
        
    prs = Presentation(str(pptx_path))
    current_date = datetime.now().strftime("%d/%m/%y")
    
    # We will update slide contents by matching slide titles/indices
    # Slide 1 (Cover): update date and footer
    # Slide 2 (Design the Offer): update slide contents
    # Slide 3 (Growth ambition and scope in China): update slide contents
    # Slide 4 (Deliver the Value): update slide contents
    # Slide 5 (Run the Company): update slide contents
    # Slide 6 (Risks, assumptions, open questions): update slide contents

    slide_updates = {
        1: {
            "UBRAN GLOW GRID": "URBAN GLOW GRID",
            "URBAN GLOW GRID": "URBAN GLOW GRID",
            "20/05/26": current_date,
            "27/05/26": current_date
        },
        2: { # Design the Offer
            "Content": (
                "Customers & problems / needs:\n"
                "\tProblem: Urban energy usage is abstract, invisible, and delayed (monthly bills). Citizens cannot optimize what they cannot see.\n"
                "\tNeed: Public spaces and municipal bodies need highly interactive, gamified tools to translate complex energy distribution grid challenges into intuitive, emotional human experiences.\n\n"
                "Solutions and concepts / architecture:\n"
                "\tThe Hero Product: An interactive digital twin & hardware layout of Paris at a strict 1:15,000 scale.\n"
                "\tSystem Architecture: A 20-zone physical wood matrix powered by an affordable ESP32 microcontroller, driven by an organic, responsive HSL \"breathing light\" logic.\n\n"
                "Value and Pricing: How we win, rough price logic:\n"
                "\tHow We Win: Gamified Social Equity. We compete on visceral engagement and extreme cost-efficiency (sub-25€ DIY building blocks vs. thousands of euros for industrial smart displays).\n"
                "\tPrice Logic: B2B/B2G subscription-based pricing for custom city layouts, paired with a freemium open-source software model for universities and schools."
            ),
            "PROJECT NAME": "URBAN GLOW GRID",
            "27/05/26": current_date
        },
        3: { # Growth ambition and scope in China
            "Content": (
                "Growth objective: What does ‘success’ look like in 3–5 years?\n"
                "\tDeploy the \"Urban Glow Grid\" platform across 15 major European and Asian Smart Cities—specifically targeting Shanghai, Shenzhen, Singapore, Paris, and Amsterdam—as an official interactive civic tech exhibition tool.\n"
                "\tReach 500,000+ online simulator users and establish commercial partnerships with municipal energy boards (e.g., EDF in France, State Grid in China).\n\n"
                "Primary target customer / segment:\n"
                "\tB2B / B2G: Smart City Museums, Science Centers, Municipal Education Boards, and Real Estate Developers seeking tangible ESG/Sustainability communication tools.\n\n"
                "Current stage:\n"
                "\tMVP / Low-cost Pilot: Low-fidelity physical 80x60cm wooden matrix prototype (25€ budget limit) paired with an active local JavaScript web simulator core."
            ),
            "PROJECT NAME": "URBAN GLOW GRID",
            "27/05/26": current_date
        },
        4: { # Deliver the Value
            "Content": (
                "Go-to-market: Channels, sales approach, type of customers:\n"
                "\tKey Message: \"Make the invisible tangible: Experience the real-time heartbeat of your city's energy grid.\"\n"
                "\tChannels: Direct inbound sales to Science & Tech museums, and partnerships with University Innovation Hubs (e.g., UTSEUS).\n"
                "\tSales Approach: High-impact, visual \"Elevator Pitch\" demonstrations utilizing the interactive dashboard to secure pilot exhibition spaces.\n\n"
                "Delivery & capacity: Where/how is it produced or hosted:\n"
                "\tHardware: Sourced entirely from open markets in mainland China (Taobao) for ultra-low-cost components; manufactured using campus fablabs.\n"
                "\tSoftware: Hosted globally via lightweight, fast-loading, client-side web hosting (GitHub Pages / Vercel), requiring zero backend database upkeep for maximum scalability.\n\n"
                "How many customers/units we can roughly serve now:\n"
                "\tCurrent Capacity: 1 physical unit (hand-wired assembly loop), but capable of supporting unlimited concurrent digital web simulator instances internationally."
            ),
            "PROJECT NAME": "URBAN GLOW GRID",
            "27/05/26": current_date
        },
        5: { # Run the Company
            "Content": (
                "People & roles: Who does what, evolution & missing skills:\n"
                "\tAbir (CEO): Product vision, narrative design, and strategic stakeholder pitch alignment.\n"
                "\tSofiane (CTO): ESP32 system engineering, HSL lighting logic programming, and hardware wire routing.\n"
                "\tMohamed (Creative Director): 3D UI/UX design, Paris map scaling, and asset creation.\n"
                "\tIsmail (Operations & Sourcing): Component procurement via Taobao, budget tracking, and fablab manufacturing coordination.\n"
                "\tNouh TARI (Head of Sales & Commercial): Driving B2B/B2G contracts, museum relationships, and market expansion.\n"
                "\tRole Evolution: As we scale, founders will transition from \"hands-on makers\" to department managers leading specialized teams.\n"
                "\tCritical Missing Skills: Industrial product designer (for scalable plastic casing), B2B Legal/Compliance Specialist, and a Customer Success Manager.\n\n"
                "Money & Company Origins: Where/how is it financed:\n"
                "\tCompany Creation: Born from a UTSEUS \"Games for Change\" academic project. The company was founded on the challenge of building a high-impact IoT tool with a strict student budget constraint.\n"
                "\tCurrent Funding: Bootstrapped via the 25€ (200 RMB) university lab grant.\n"
                "\tScale Financing Plan: Transitioning to regional green incubation grants and academic innovation funds in Shanghai/Europe to fund V2 manufacturing.\n\n"
                "Legal & governance:\n"
                "\tCompliance with open-source software distributions (MIT License framework).\n"
                "\tAdherence to municipal lab health, safety, and electrical standards."
            ),
            "PROJECT NAME": "URBAN GLOW GRID",
            "27/05/26": current_date
        },
        6: { # Risks, assumptions, open questions
            "Content": (
                "Key assumptions: Things we take for granted:\n"
                "\tInstitutional buyers prefer physical, tactile, wooden layouts over plain display interfaces.\n"
                "\tReal-time carbon intensity data metrics (RTE/ADEME API models) can be clearly simplified for everyday urban citizens.\n\n"
                "Unknowns / must-learn (Must be validated before big investment & Action Plan):\n"
                "\tWill users engage long-term or treat it as a short-lived gimmick? -> Action Plan: Deploy a 4-week beta test in a Shanghai museum and track average session duration and repeat visits using hidden software analytics.\n"
                "\tWhat is the exact wear-and-tear lifespan of cheap mechanical switches? -> Action Plan: Build an automated mechanical testing rig in the lab to simulate 10,000 presses and record failure rates before mass ordering.\n"
                "\tWhat are the exact certification costs and compliance pathways for educational hardware? -> Action Plan: Consult a legal expert specializing in EU CE marking and China CCC certification by Q3 to map out regulatory costs.\n\n"
                "Major risks: What could break our plan:\n"
                "\tSupply Chain Disruption: Fluctuations in components or microchip availability affecting production costs.\n"
                "\tThe Novelty Fatigue: If the interactive experience feels too static, institutions will cancel their recurring software subscriptions."
            ),
            "PROJECT NAME": "URBAN GLOW GRID",
            "27/05/26": current_date
        }
    }

    # Iterate slides and shapes to perform updates
    for i, slide in enumerate(prs.slides, start=1):
        updates = slide_updates.get(i, {})
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                orig_text = shape.text.strip()
                # Check direct text matching or specific placeholder types
                if "Content Placeholder" in shape.name and "Content" in updates:
                    shape.text = updates["Content"]
                elif "Footer Placeholder" in shape.name and "PROJECT NAME" in updates:
                    shape.text = updates["PROJECT NAME"]
                elif "Footer Placeholder" in shape.name and "UBRAN GLOW GRID" in orig_text:
                    shape.text = updates["UBRAN GLOW GRID"]
                elif "Footer Placeholder" in shape.name and "URBAN GLOW GRID" in orig_text:
                    shape.text = "URBAN GLOW GRID"
                elif "Date Placeholder" in shape.name and ("20/05/26" in orig_text or "27/05/26" in orig_text):
                    shape.text = current_date
                elif orig_text in ["20/05/26", "27/05/26"]:
                    shape.text = current_date
                elif "UBRAN" in orig_text:
                    shape.text = orig_text.replace("UBRAN", "URBAN")
                    
    # Save the updated presentation
    output_path = pptx_path.parent / "Three‑Engine Scale‑Up Canvas-UrbanGlowGrid_v1.5.pptx"
    prs.save(str(output_path))
    prs.save(str(pptx_path))
    print(f"SUCCESS: PowerPoint Canvas updated and overwritten {pptx_path.name}")

if __name__ == "__main__":
    main()
